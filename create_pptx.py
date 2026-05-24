#!/usr/bin/env python3
"""
CarbonTrack — editable PPTX generator.
All 15 slides, RTL-correct layout, sizes derived from presentation.html (1600×900).

Run:  pip install python-pptx Pillow requests
      python create_pptx.py
"""
import io, sys, subprocess, requests
from pathlib import Path
from PIL import Image

try:
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    subprocess.run([sys.executable, '-m', 'pip', 'install', 'python-pptx', '-q'], check=True)
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

BASE = Path(__file__).parent
IMGS = BASE / 'images'
def imgp(name): return str(IMGS / name)
def has(name):  return (IMGS / name).exists()

# ── Colours ───────────────────────────────────────────────────────────────────
CREAM  = RGBColor(0xF5, 0xF2, 0xEB)
G9     = RGBColor(0x1B, 0x43, 0x32)
G5     = RGBColor(0x52, 0xB7, 0x88)
G3     = RGBColor(0x95, 0xD5, 0xB2)
G1     = RGBColor(0xD8, 0xF3, 0xE3)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0x8A, 0xAA, 0x8E)
SUB    = RGBColor(0x4A, 0x7C, 0x59)
RED    = RGBColor(0xC0, 0x39, 0x2B)
BORDER = RGBColor(0xC8, 0xDD, 0xC0)
BLU_BG = RGBColor(0xC8, 0xE6, 0xFA)
BLU_FG = RGBColor(0x0D, 0x3B, 0x6E)
ORG_BG = RGBColor(0xFD, 0xE8, 0xB8)
ORG_FG = RGBColor(0x6B, 0x3A, 0x0D)
GLD_BG = RGBColor(0xFE, 0xF3, 0xC7)
GLD_FG = RGBColor(0x92, 0x40, 0x0E)
ORANGE = RGBColor(0xC0, 0x7A, 0x2A)

# ── Presentation (16:9, matching 1600×900) ────────────────────────────────────
prs = Presentation()
prs.slide_width  = Emu(9144000)   # 25.4 cm  (standard 16:9)
prs.slide_height = Emu(5143500)   # 14.29 cm
W = prs.slide_width
H = prs.slide_height
BL = prs.slide_layouts[6]

# ── Unit helpers: HTML px → PPTX EMU / pt ─────────────────────────────────────
def px(v): return int(W * v / 1600)   # horizontal pixel
def py(v): return int(H * v / 900)    # vertical pixel
def fz(v): return round(v * 0.75)     # CSS px → font pt

# ── RTL column layout helpers ─────────────────────────────────────────────────
PAD = 48   # slide padding in CSS px

def col_w(n, pad=PAD, gap=16):
    """Width (px-equivalent EMU) of each equal column in RTL row of n."""
    return (W - px(2*pad) - px(gap)*(n-1)) // n

def col_x(i, n, pad=PAD, gap=16):
    """Left-edge x of item i (i=0 = RIGHTMOST) in RTL row of n equal columns."""
    cw = col_w(n, pad, gap)
    return W - px(pad) - (i+1)*cw - i*px(gap)

# ── Low-level helpers ──────────────────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(BL)

def bg(sl, color):
    f = sl.background.fill; f.solid(); f.fore_color.rgb = color

def box(sl, lx, ty, rw, rh, fill=None, line=None, lw=Pt(1.5)):
    shp = sl.shapes.add_shape(1, lx, ty, rw, rh)
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:    shp.fill.background()
    if line: shp.line.color.rgb = line; shp.line.width = lw
    else:    shp.line.fill.background()
    return shp

def circle(sl, cx, cy, d, fill=G5, line=None, line_color=None):
    shp = sl.shapes.add_shape(9, cx-d//2, cy-d//2, d, d)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line and line_color:
        shp.line.color.rgb = line_color; shp.line.width = Pt(line)
    else:
        shp.line.fill.background()
    return shp

PPTX_OK = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.wmf'}

def add_pic(sl, path, lx, ty, rw, rh=None):
    try:
        ext = Path(path).suffix.lower() if isinstance(path, str) else ''
        if ext not in PPTX_OK:
            buf = io.BytesIO()
            Image.open(path).convert('RGBA').save(buf, 'PNG'); buf.seek(0)
            path = buf
        if rh: sl.shapes.add_picture(path, lx, ty, rw, rh)
        else:  sl.shapes.add_picture(path, lx, ty, rw)
    except Exception as e:
        print(f'  [!] pic: {e}')

def txb(sl, text, lx, ty, rw, rh,
        size=18, bold=False, color=G9, italic=False,
        align=PP_ALIGN.RIGHT, rtl=True):
    bx = sl.shapes.add_textbox(lx, ty, rw, rh)
    tf = bx.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    rn = p.add_run(); rn.text = text
    rn.font.size = Pt(size); rn.font.bold = bold
    rn.font.italic = italic; rn.font.color.rgb = color
    rn.font.name = 'Calibri'
    if rtl: p._p.get_or_add_pPr().set('rtl', '1')
    return bx

def txb_lines(sl, lines, lx, ty, rw, rh,
              size=18, bold=False, color=G9,
              align=PP_ALIGN.RIGHT, rtl=True):
    bx = sl.shapes.add_textbox(lx, ty, rw, rh)
    tf = bx.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.alignment = align
        rn = p.add_run(); rn.text = line
        rn.font.size = Pt(size); rn.font.bold = bold
        rn.font.color.rgb = color; rn.font.name = 'Calibri'
        if rtl: p._p.get_or_add_pPr().set('rtl', '1')
    return bx

def lbl(sl, text, ty=44):
    """Section label — 84px, G5, bold, right-aligned."""
    # Green accent bar (mimics .label::before) on the right edge
    bar_h = py(5)
    box(sl, W - px(PAD) - px(28), py(ty) + py(35), px(28), bar_h, fill=G5)
    txb(sl, text, px(PAD), py(ty), W - px(2*PAD) - px(38), py(84),
        size=fz(84), bold=True, color=G5)

def title(sl, text, ty=140):
    """Slide h2 title — 52px, G9, bold."""
    txb(sl, text, px(PAD), py(ty), W - px(2*PAD), py(90),
        size=fz(52), bold=True, color=G9)


# ══════════════════════════════════════════════════════════════════════════════
# S1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, CREAM)
logo_sz = px(280)
if has('logo.png'):
    add_pic(sl, imgp('logo.png'), W//2 - logo_sz//2, py(120), logo_sz, logo_sz)

txb(sl, 'CarbonTrack', px(PAD), py(430), W - px(2*PAD), py(160),
    size=fz(130), bold=True, color=G9, align=PP_ALIGN.CENTER)

# Partners bar — in HTML: flex row centered (not RTL-flipped because justify-content:center)
p_imgs = [
    ('NetiviIsrael.svg.png',              'נתיבי ישראל'),
    ('Google_2015_logo.svg.png',          'Google'),
    ('Israel_National_Digital_Agency_logo.png', 'רשות הדיגיטל'),
    ('sti-logo.webp',                     'STI'),
    ('yazamut_logo_white.webp',           'יזמות'),
]
pw = (W - px(2*PAD) - px(56)*4) // len(p_imgs)
pcx = px(PAD)
for pname, palt in p_imgs:
    if has(pname):
        add_pic(sl, imgp(pname), pcx, H - py(140), pw, py(90))
    else:
        txb(sl, palt, pcx, H - py(130), pw, py(80), size=12, color=SUB, align=PP_ALIGN.CENTER)
    pcx += pw + px(56)
print('S1 Title')


# ══════════════════════════════════════════════════════════════════════════════
# S2 — HOOK  (centered layout — not RTL-dependent)
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, CREAM)
# Three icons — centered in equal thirds (HTML uses centered flex, not RTL reorder)
icons3 = [('🏗️', 'חומרי בנייה'), ('☁️', 'פליטות CO₂'), ('🌍', 'פגיעה באטמוספירה')]
cw3 = W // 3
for i, (icon, cap) in enumerate(icons3):
    cx3 = i * cw3
    txb(sl, icon, cx3, py(130), cw3, py(160), size=fz(100), align=PP_ALIGN.CENTER, rtl=False)
    txb(sl, cap,  cx3, py(300), cw3, py(80),  size=fz(46),  bold=True, color=G9, align=PP_ALIGN.CENTER)

# Arrows ← between the three icons
txb(sl, '←', cw3 - px(60),   py(160), px(120), py(100), size=fz(72), color=G5, align=PP_ALIGN.CENTER, rtl=False)
txb(sl, '←', 2*cw3 - px(60), py(160), px(120), py(100), size=fz(72), color=G5, align=PP_ALIGN.CENTER, rtl=False)

box(sl, W//2 - px(70), py(440), px(140), py(4), fill=BORDER)

txb_lines(sl, ['אי אפשר לצמצם', 'את מה שלא מודדים'],
          px(PAD), py(470), W - px(2*PAD), py(380),
          size=fz(100), bold=True, color=G9, align=PP_ALIGN.CENTER)
print('S2 Hook')


# ══════════════════════════════════════════════════════════════════════════════
# S3 — TEAM  (RTL grid: אסף=rightmost, שירה=leftmost)
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, CREAM)
lbl(sl, 'הצוות')
txb(sl, 'סטודנטים להנדסת חשמל ומחשבים',
    px(PAD), py(115), W - px(2*PAD), py(55), size=fz(32), color=SUB)

members = [
    ('asaf.jpg',  'א', 'אסף שדה',    'לוחם יחידת אגוז\nמשרת מילואים פעיל'),
    ('gal.jpg',   'ג', 'גל הראל',    'יחידה 9900\nמשרתת מילואים פעילה'),
    ('mika.jpg',  'מ', 'מיקה הרפז',  'יחידה 8200 · מודיעין אותות\nמשרתת מילואים פעילה'),
    ('shira.jpg', 'ש', 'שירה זיגדון', 'יחידה 8200 · חוקרת RF\nמשרתת מילואים פעילה'),
]
card_t = py(178)
card_h = H - card_t - py(16)
photo_h = int(card_h * 0.62)
# RTL: i=0 (אסף) = rightmost
for i, (pfile, ini, name, unit) in enumerate(members):
    cx4 = col_x(i, 4, gap=20)
    cw4 = col_w(4, gap=20)
    box(sl, cx4, card_t, cw4, card_h, fill=WHITE, line=BORDER)
    if has(pfile):
        add_pic(sl, imgp(pfile), cx4+px(4), card_t+py(4), cw4-px(8), photo_h)
    else:
        box(sl, cx4+px(4), card_t+py(4), cw4-px(8), photo_h, fill=G1)
        txb(sl, ini, cx4, card_t+photo_h//4, cw4, photo_h//2,
            size=fz(52), bold=True, color=G5, align=PP_ALIGN.CENTER)
    info_y = card_t + photo_h + py(14)
    txb(sl, name, cx4+px(18), info_y, cw4-px(36), py(55), size=fz(42), bold=True, color=G9)
    txb_lines(sl, unit.split('\n'),
              cx4+px(18), info_y+py(60), cw4-px(36), py(90), size=fz(28), color=MUTED)
print('S3 Team')


# ══════════════════════════════════════════════════════════════════════════════
# S4 — PROBLEM  (RTL: fact1=right, fact2=left)
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, CREAM)
# Label positioned top-right as in HTML (.problem-label { position:absolute; top:44px; right:60px })
txb(sl, 'הבעיה', W - px(60) - px(300), py(44), px(300), py(100),
    size=fz(84), bold=True, color=G5)

txb(sl, '81M טון', px(PAD), py(95), W - px(2*PAD), py(240),
    size=fz(148), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txb(sl, 'CO₂ נפלטים מדי שנה בישראל',
    px(PAD), py(335), W - px(2*PAD), py(70), size=fz(46), bold=True, color=SUB, align=PP_ALIGN.CENTER)

# Two fact cards — RTL: first fact (שינויים) = RIGHT, second (2500) = LEFT
fw = col_w(2, gap=20); fy = py(415); fh = py(165)
# Fact 1 (rightmost, i=0)
fx1 = col_x(0, 2, gap=20)
box(sl, fx1, fy, fw, fh, fill=WHITE, line=BORDER, lw=Pt(2))
txb(sl, 'שינויים קיצוניים במזג האוויר', fx1+px(20), fy+py(18), fw-px(30), py(70),
    size=fz(36), bold=True, color=G9)
txb(sl, 'גלי חום, הצפות ונזקים של מיליונים בכל חורף',
    fx1+px(20), fy+py(92), fw-px(30), py(65), size=fz(28), color=MUTED)

# Fact 2 (leftmost, i=1)
fx2 = col_x(1, 2, gap=20)
box(sl, fx2, fy, fw, fh, fill=WHITE, line=BORDER, lw=Pt(2))
txb(sl, '2,500', fx2+px(20), fy+py(14), fw-px(30), py(80), size=fz(42), bold=True, color=G9)
txb(sl, 'ישראלים מתים מדי שנה כתוצאה ישירה מזיהום אוויר',
    fx2+px(20), fy+py(98), fw-px(30), py(65), size=fz(28), color=MUTED)

# Infra dark banner (full width)
bx = px(PAD); by = py(598); bw = W - px(2*PAD); bh = py(272)
box(sl, bx, by, bw, bh, fill=G9)
txb(sl, '🏗️', bx+px(30), by+py(20), px(120), py(130), size=fz(90), rtl=False)
txb(sl, 'יותר מ-10%', bx+px(160), by+py(15), bw-px(180), py(115),
    size=fz(80), bold=True, color=G3)
txb(sl, 'מסך הפליטות הישראליות מגיע מענף התשתיות בלבד',
    bx+px(160), by+py(133), bw-px(180), py(70), size=fz(32), color=WHITE)
txb(sl, '← זה הנתח שאנחנו באים לטפל בו',
    bx+px(160), by+py(205), bw-px(180), py(55), size=fz(26), bold=True, color=G3)
print('S4 Problem')


# ══════════════════════════════════════════════════════════════════════════════
# S5 — GHG PROTOCOL
# RTL layout: Scope1+2 group = RIGHT half, Scope3 = LEFT half
# Within Scope1+2: Scope1 = rightmost, Scope2 = to its left
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, CREAM)
lbl(sl, 'המדידה')
title(sl, 'פרוטוקול GHG — כיצד מודדים?')

sy = py(210); sh = H - sy - py(16)
grp_gap = px(16)
inner_gap = px(16)
avail = W - px(2*PAD) - grp_gap
grp_w = avail // 2          # each main group width
sc_w  = (grp_w - inner_gap) // 2   # each individual scope width (in the S1+2 group)
bottom_lbl_h = py(60)
scope_h = sh - bottom_lbl_h - py(14)  # height of scope card itself

# LEFT side: Scope 3 (second child in HTML → leftmost in RTL)
s3x = px(PAD)
box(sl, s3x, sy, grp_w, scope_h, fill=G9)
txb(sl, 'SCOPE 3 — הבעיה האמיתית',
    s3x+px(18), sy+py(18), grp_w-px(30), py(55), size=fz(20), bold=True, color=G3, rtl=False)
txb(sl, 'שרשרת האספקה',
    s3x+px(18), sy+py(80), grp_w-px(30), py(70), size=fz(44), bold=True, color=WHITE)
txb_lines(sl,
    ['ייצור חומרים, הובלה, קבלני משנה —', 'וזה בדיוק מה שאנחנו פותרים.'],
    s3x+px(18), sy+py(158), grp_w-px(30), py(120), size=fz(28), color=WHITE)
txb(sl, '70%–90% מסך הפליטות',
    s3x+px(18), sy+py(295), grp_w-px(30), py(65), size=fz(32), bold=True, color=G3)
# Bar
box(sl, s3x+px(18), sy+py(368), grp_w-px(36), py(10), fill=RGBColor(0x40,0x70,0x58))
box(sl, s3x+px(18), sy+py(368), int((grp_w-px(36))*0.85), py(10), fill=G3)
# Bottom ✗ label
txb(sl, '✗ קשה למדידה',
    s3x, sy+scope_h+py(14), grp_w, bottom_lbl_h,
    size=fz(34), bold=True, color=G9, align=PP_ALIGN.CENTER)

# RIGHT side: Scope 1+2 group (first child in HTML → rightmost in RTL)
s12x = px(PAD) + grp_w + grp_gap  # left edge of the Scope1+2 group

# Scope 2 (second sub-child → leftmost of the group, left of Scope 1)
s2x = s12x
box(sl, s2x, sy, sc_w, scope_h, fill=ORG_BG)
txb(sl, 'SCOPE 2', s2x+px(16), sy+py(18), sc_w-px(24), py(42), size=fz(20), bold=True, color=ORG_FG, rtl=False)
txb(sl, 'אנרגיה נרכשת', s2x+px(16), sy+py(68), sc_w-px(24), py(65), size=fz(38), bold=True, color=ORG_FG)
txb(sl, 'חשמל וחום נרכש מספקים.',
    s2x+px(16), sy+py(145), sc_w-px(24), py(130), size=fz(28), color=ORG_FG)

# Scope 1 (first sub-child → rightmost of the group)
s1x = s12x + sc_w + inner_gap
box(sl, s1x, sy, sc_w, scope_h, fill=BLU_BG)
txb(sl, 'SCOPE 1', s1x+px(16), sy+py(18), sc_w-px(24), py(42), size=fz(20), bold=True, color=BLU_FG, rtl=False)
txb(sl, 'פליטות ישירות', s1x+px(16), sy+py(68), sc_w-px(24), py(65), size=fz(38), bold=True, color=BLU_FG)
txb(sl, 'שריפת דלק ישירה בציוד הארגון.',
    s1x+px(16), sy+py(145), sc_w-px(24), py(130), size=fz(28), color=BLU_FG)

# Bottom ✓ label spanning the full Scope1+2 group
txb(sl, '✓ קל למדידה',
    s12x, sy+scope_h+py(14), grp_w, bottom_lbl_h,
    size=fz(34), bold=True, color=G9, align=PP_ALIGN.CENTER)
print('S5 GHG')


# ══════════════════════════════════════════════════════════════════════════════
# S6 — GLOBAL CONTEXT
# Paris banner spans full width; then 3 cards in RTL: 2030=right, 2050=mid, היום=left
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, CREAM)
lbl(sl, 'ההקשר העולמי')
title(sl, 'ישראל מחויבת — חסר הכלי')

# Paris full-width banner (grid-column: span 3)
par_y = py(208); par_h = py(105)
box(sl, px(PAD), par_y, W - px(2*PAD), par_h, fill=G9)
txb(sl, '🕊️   הסכם פריז 2015',
    px(PAD), par_y, W - px(2*PAD), par_h,
    size=fz(42), bold=True, color=G3, align=PP_ALIGN.CENTER)

# Three cards below — RTL grid order: 2030 (rightmost), 2050, היום (leftmost)
cards = [
    ('2030', 'ישראל: צמצום פליטות בכ-30%\nביחס ל-2015',             WHITE,  G5,     SUB),
    ('2050', 'יעד Net Zero —\nניטרליות פחמנית מלאה',                WHITE,  G5,     SUB),
    ('היום', 'אין כלי מדידה מעשי —\nCarbonTrack ממלאת את הפער',     GLD_BG, GLD_FG, RGBColor(0x78,0x35,0x0F)),
]
cy = par_y + par_h + py(16); ch = H - cy - py(16)
for i, (yr, body, cbg, yrc, bc) in enumerate(cards):
    cx5 = col_x(i, 3)           # i=0 → rightmost (2030)
    cw5 = col_w(3)
    box(sl, cx5, cy, cw5, ch, fill=cbg, line=BORDER)
    txb(sl, yr, cx5+px(20), cy+py(20), cw5-px(30), py(185), size=fz(72), bold=True, color=yrc)
    txb_lines(sl, body.split('\n'),
              cx5+px(20), cy+py(215), cw5-px(30), ch-py(230), size=fz(28), color=bc)
print('S6 Global')


# ══════════════════════════════════════════════════════════════════════════════
# S7 — BEFORE / AFTER
# RTL flex row: .before (first) = RIGHT half, .after (second) = LEFT half
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, WHITE)
hw = W // 2

# LEFT half: dark green ".after" (second child in HTML)
box(sl, 0, 0, hw, H, fill=G9)
txb(sl, '✅  עם CarbonTrack', px(60), py(130), hw-px(80), py(55),
    size=fz(34), bold=True, color=G3, rtl=False)
txb(sl, 'העתיד הקרוב', px(60), py(198), hw-px(80), py(80), size=fz(62), bold=True, color=WHITE)
benefits = [
    'Scope 3 מלא לפי GHG Protocol',
    'תוצאות בזמן אמת',
    'חישוב מבוסס AI, לא הצהרה בלבד',
    'מעקב רגולטורי מפוקח',
]
for i, b in enumerate(benefits):
    txb(sl, f'✓  {b}', px(60), py(315) + i*py(110), hw-px(80), py(95), size=fz(32), color=WHITE)

# RIGHT half: white ".before" (first child in HTML)
box(sl, hw, 0, hw, H, fill=WHITE)
txb(sl, '❌  ללא CarbonTrack', hw+px(60), py(130), hw-px(80), py(55),
    size=fz(34), bold=True, color=RED, rtl=False)
txb(sl, 'המצב היום', hw+px(60), py(198), hw-px(80), py(80), size=fz(62), bold=True, color=G9)
txb(sl, '0', hw+px(60), py(270), hw-px(80), py(360),
    size=fz(200), bold=True, color=RGBColor(0xE0,0x7B,0x7B), align=PP_ALIGN.CENTER)
txb(sl, 'חברות תשתית ישראליות מחשבות Scope 3',
    hw+px(60), py(660), hw-px(80), py(65), size=fz(30), bold=True, color=SUB, align=PP_ALIGN.CENTER)
txb(sl, '✗  חברות בינלאומיות — לא מותאמות לשוק התשתיות הישראלי',
    hw+px(60), py(740), hw-px(80), py(90), size=fz(28), color=SUB)

# VS badge on the dividing line
vs_w, vs_h = px(80), py(60)
box(sl, hw - vs_w//2, H//2 - vs_h//2, vs_w, vs_h, fill=CREAM, line=BORDER)
txb(sl, 'VS', hw - vs_w//2, H//2 - vs_h//2, vs_w, vs_h,
    size=fz(18), bold=True, color=G9, align=PP_ALIGN.CENTER, rtl=False)
print('S7 Before/After')


# ══════════════════════════════════════════════════════════════════════════════
# S8 — VIDEO DEMO  (placeholder — video can't be embedded easily)
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, CREAM)
lbl(sl, 'הדגמה')
title(sl, 'ראו איך זה עובד')

# Video frame placeholder
vx = px(PAD); vy = py(220)
vw = W - px(2*PAD); vh = H - vy - py(30)
box(sl, vx, vy, vw, vh, fill=G9, line=BORDER)
txb(sl, '▶', vx, vy + vh//2 - py(100), vw, py(160),
    size=120, bold=False, color=G3, align=PP_ALIGN.CENTER, rtl=False)
txb(sl, 'carbontrack_cut.mp4', vx, vy + vh//2 + py(60), vw, py(70),
    size=fz(28), color=G3, align=PP_ALIGN.CENTER, rtl=False)
txb(sl, '(פתח את קובץ הווידאו ידנית למצגת)', vx, vy + vh//2 + py(135), vw, py(55),
    size=fz(22), color=MUTED, align=PP_ALIGN.CENTER)
print('S8 Video')


# ══════════════════════════════════════════════════════════════════════════════
# S9 — WHAT-IF  (full-slide image + label)
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, CREAM)
lbl(sl, 'סימולטור What-If', ty=16)
if has('WHAT-IF.png'):
    add_pic(sl, imgp('WHAT-IF.png'), px(20), py(95), W - px(40), H - py(110))
else:
    txb(sl, '[תמונת What-If]', 0, H//2, W, py(80),
        size=32, color=MUTED, align=PP_ALIGN.CENTER)
print('S9 What-If')


# ══════════════════════════════════════════════════════════════════════════════
# S10 — AI HELPER  (full-slide image + label)
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, CREAM)
lbl(sl, 'עוזר AI', ty=16)
if has('AI.png'):
    add_pic(sl, imgp('AI.png'), px(20), py(95), W - px(40), H - py(110))
else:
    txb(sl, '[תמונת עוזר AI]', 0, H//2, W, py(80),
        size=32, color=MUTED, align=PP_ALIGN.CENTER)
print('S10 AI')


# ══════════════════════════════════════════════════════════════════════════════
# S11 — HOW IT WORKS (3 steps)
# RTL flex row: Step1 (כתב כמויות) = RIGHTMOST, Step3 (דאשבורד) = LEFTMOST
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, CREAM)
lbl(sl, 'הפתרון')
title(sl, 'איך זה עובד?')

# 3 equal squares with 2 arrows between them
arr_w = px(44)
sq_w  = (W - px(2*PAD) - 2*arr_w) // 3
sq_y  = py(208); sq_h = sq_w   # square (aspect-ratio:1/1)

steps = [
    # (pfile,        icon,  title,              desc,                                    bg,    fg)
    ('papers.jpg',  '📄', 'העלאת כתב כמויות', 'מסמך המפרט את כלל העבודות\nוהחומרים הנדרשים למיזם.', WHITE, G9),
    (None,          '🔢', 'חישוב פליטה',       'ע"י מקדמי המרה בינלאומיים\n✓ GHG Protocol Approved',  G9,   WHITE),
    (None,          '📊', 'דאשבורד נתונים',    'ניתוח חי לפי פרויקט,\nחומר וקבלן',                    WHITE, G9),
]
# RTL: step 0 = rightmost, step 2 = leftmost
for i, (pfile, icon, sname, sdesc, sbg, sfg) in enumerate(steps):
    # col_x with i=0 rightmost
    sx = col_x(i, 3, gap=arr_w // px(1))   # gap in "px" units
    # recompute manually for clarity
    sx = W - px(PAD) - (i+1)*sq_w - i*arr_w

    box(sl, sx, sq_y, sq_w, sq_h,
        fill=sbg, line=BORDER if sbg != G9 else None, lw=Pt(2))
    if pfile and has(pfile):
        add_pic(sl, imgp(pfile), sx+px(28), sq_y+px(32), sq_w-px(56), py(120))
    else:
        txb(sl, icon, sx, sq_y+px(28), sq_w, py(145), size=fz(80), align=PP_ALIGN.CENTER, rtl=False)
    txb(sl, sname, sx+px(20), sq_y+py(220), sq_w-px(30), py(80),
        size=fz(38), bold=True, color=sfg, align=PP_ALIGN.CENTER)
    txb_lines(sl, sdesc.split('\n'),
              sx+px(20), sq_y+py(310), sq_w-px(30), sq_h-py(330),
              size=fz(26), color=sfg if sfg==WHITE else MUTED, align=PP_ALIGN.CENTER)
    # Arrow to the LEFT of each step (except last/leftmost)
    if i < 2:
        ax = sx - arr_w
        txb(sl, '←', ax, sq_y + sq_h//2 - py(40), arr_w, py(80),
            size=fz(44), color=G5, align=PP_ALIGN.CENTER, rtl=False)
print('S11 How it works')


# ══════════════════════════════════════════════════════════════════════════════
# S12 — TECH FLOW
# RTL: 📄 כתב כמויות = RIGHTMOST, 📊 דוח BI = LEFTMOST
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, CREAM)
lbl(sl, 'הטכנולוגיה')
title(sl, 'מכתב כמויות לדוח פליטות — בדקות')

flow = [
    ('📄', 'כתב כמויות',  'Excel / CSV'),
    ('🧠', 'זיהוי AI',    'NLP הנדסי'),
    ('⚖️', 'המרת משקל',  'קטלוג חומרים'),
    ('🔢', 'חישוב פליטה', 'ושמירתה ב-DB'),
    ('📊', 'דוח BI',      'דאשבורד חי'),
]
n_steps = len(flow)
arr_w2  = px(30)
fl_w    = (W - px(2*PAD) - arr_w2*(n_steps-1)) // n_steps
fl_y    = py(215); fl_h = py(280)

for i, (icon, ftitle, fsub) in enumerate(flow):
    # i=0 = rightmost
    fx = W - px(PAD) - (i+1)*fl_w - i*arr_w2
    box(sl, fx, fl_y, fl_w, fl_h, fill=WHITE, line=BORDER)
    txb(sl, icon,   fx, fl_y+py(18), fl_w, py(110), size=fz(52), align=PP_ALIGN.CENTER, rtl=False)
    txb(sl, ftitle, fx+px(8), fl_y+py(135), fl_w-px(16), py(75), size=fz(28), bold=True, color=G9, align=PP_ALIGN.CENTER)
    txb(sl, fsub,   fx+px(8), fl_y+py(212), fl_w-px(16), py(58), size=fz(22), color=MUTED, align=PP_ALIGN.CENTER)
    if i < n_steps-1:
        txb(sl, '←', fx - arr_w2, fl_y + fl_h//2 - py(30), arr_w2, py(60),
            size=fz(24), color=G5, align=PP_ALIGN.CENTER, rtl=False)

# Notes — RTL: Self-Learning = RIGHT, ביצועים = LEFT
notes = [
    ('📈 Self-Learning', 'משתפרת בכל קובץ שנכנס'),
    ('⚡ ביצועים',       'אלפי שורות בדקות, progress בזמן אמת'),
]
ny = py(515); nh = H - ny - py(16)
for i, (nt, nd) in enumerate(notes):
    nx = col_x(i, 2, gap=14)   # i=0 = rightmost
    nw = col_w(2, gap=14)
    box(sl, nx, ny, nw, nh, fill=WHITE, line=BORDER)
    txb(sl, nt, nx+px(24), ny+py(22), nw-px(40), py(65), size=fz(28), bold=True, color=RGBColor(0x2D,0x6A,0x4F))
    txb(sl, nd, nx+px(24), ny+py(92), nw-px(40), nh-py(108), size=fz(22), color=MUTED)
print('S12 Tech flow')


# ══════════════════════════════════════════════════════════════════════════════
# S13 — ROADMAP
# RTL: PoC ראשוני = RIGHTMOST (i=0), חברות תשתית = LEFTMOST (i=5)
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, CREAM)
lbl(sl, 'מפת דרכים')
title(sl, 'מה עשינו ולאן אנחנו הולכים')

milestones = [
    # (date,        title,                    icon,  status,    up?)
    ('ינו׳ 2026', 'PoC ראשוני',             '✅', 'done',    True),
    ('מרץ 2026',  'AI Integration',          '✅', 'done',    False),
    ('אפר׳ 2026', 'Dashboard MVP',           '✅', 'done',    True),
    ('מאי 2026',  'Pilot נתיבי ישראל',       '🔧', 'current', False),
    ('2026',      'הטמעה בנתיבי ישראל',     '⚙️', 'future',  True),
    ('2027',      'חברות תשתית נוספות',     '🏗️', 'future',  False),
]
n_ms  = len(milestones)
line_y = py(495)
lx0, lxN = px(80), W - px(80)
box(sl, lx0, line_y - py(4), lxN - lx0, py(8), fill=G5)  # timeline bar

ms_cw = (lxN - lx0) // n_ms
dot_d  = py(64)

for i, (date, mtitle, icon, status, up) in enumerate(milestones):
    # i=0 = rightmost
    cx6 = W - px(80) - ms_cw//2 - i*ms_cw
    dot_color = G5 if status == 'current' else (G9 if status == 'done' else MUTED)
    circle(sl, cx6, line_y, dot_d, fill=dot_color)
    txb(sl, icon, cx6 - dot_d//2, line_y - dot_d//2, dot_d, dot_d,
        size=fz(28), align=PP_ALIGN.CENTER, rtl=False)
    tw = ms_cw - px(10); tx = cx6 - tw//2
    if up:
        txb(sl, date,   tx, py(320), tw, py(48), size=fz(22), color=G5,  align=PP_ALIGN.CENTER)
        txb(sl, mtitle, tx, py(368), tw, py(95), size=fz(26), bold=True, color=G9, align=PP_ALIGN.CENTER)
    else:
        if status == 'current':
            # "📍 היום" badge above the dot
            bw2 = px(130); bh2 = py(44)
            box(sl, cx6-bw2//2, line_y - dot_d//2 - bh2 - py(8), bw2, bh2, fill=G5)
            txb(sl, '📍 היום', cx6-bw2//2, line_y - dot_d//2 - bh2 - py(6),
                bw2, bh2, size=fz(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txb(sl, date,   tx, py(572), tw, py(48), size=fz(22), color=G5,  align=PP_ALIGN.CENTER)
        txb(sl, mtitle, tx, py(620), tw, py(95), size=fz(26), bold=True, color=G9, align=PP_ALIGN.CENTER)
        if i == n_ms - 1:
            txb(sl, 'רכבת ישראל, מקורות ועוד',
                tx, py(718), tw, py(50), size=fz(20), color=MUTED, align=PP_ALIGN.CENTER)
print('S13 Roadmap')


# ══════════════════════════════════════════════════════════════════════════════
# S14 — CONVERGENCE (3 circles + logo)
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, CREAM)

# 3 overlapping circles in a triangle formation, centred on slide
# HTML: 640×580 container centred in 1600×900 slide
cont_w = px(640); cont_h = py(580)
cont_x = W//2 - cont_w//2; cont_y = H//2 - cont_h//2
circ_d = px(260)

def conv_circle(cx, cy, fill_hex, border_hex, emoji):
    r, g, b = int(fill_hex[0:2],16), int(fill_hex[2:4],16), int(fill_hex[4:6],16)
    rb, gb, bb = int(border_hex[0:2],16), int(border_hex[2:4],16), int(border_hex[4:6],16)
    shp = sl.shapes.add_shape(9, cx - circ_d//2, cy - circ_d//2, circ_d, circ_d)
    shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor(r,g,b)
    shp.line.color.rgb = RGBColor(rb,gb,bb); shp.line.width = Pt(3)
    txb(sl, emoji, cx - circ_d//2, cy - circ_d//2 + py(55), circ_d, circ_d - py(55),
        size=fz(100), align=PP_ALIGN.CENTER, rtl=False)

# top-right: 💰
conv_circle(cont_x + cont_w - px(10) - circ_d//2,
            cont_y + py(10) + circ_d//2,
            '21a85321', '34a853', '💰')
# top-left: 🏛️
conv_circle(cont_x + px(10) + circ_d//2,
            cont_y + py(10) + circ_d//2,
            '1b43321f', '1b4332', '🏛️')
# bottom-center: 🤖
conv_circle(cont_x + px(190) + circ_d//2,
            cont_y + cont_h - py(10) - circ_d//2,
            '52b7881f', '52b788', '🤖')

# Logo circle in the centre
logo_d = px(150)
logo_cx = cont_x + cont_w//2; logo_cy = cont_y + int(cont_h * 0.43)
shp_logo = sl.shapes.add_shape(9,
    logo_cx - logo_d//2, logo_cy - logo_d//2, logo_d, logo_d)
shp_logo.fill.solid(); shp_logo.fill.fore_color.rgb = WHITE
shp_logo.line.color.rgb = G3; shp_logo.line.width = Pt(3)
if has('logo.png'):
    add_pic(sl, imgp('logo.png'),
            logo_cx - logo_d//2, logo_cy - logo_d//2, logo_d, logo_d)
print('S14 Convergence')


# ══════════════════════════════════════════════════════════════════════════════
# S15 — CLOSING
# RTL flex row: text (first child) = RIGHT, QR card (second child) = LEFT
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, CREAM)

# QR card — LEFT side (second child in RTL)
qr_card_w = px(390); qr_card_h = H - py(80)
qr_x = px(PAD); qr_y = py(40)
box(sl, qr_x, qr_y, qr_card_w, qr_card_h, fill=WHITE, line=BORDER, lw=Pt(2))
txb(sl, 'SCAN ME', qr_x, qr_y+py(20), qr_card_w, py(40),
    size=fz(13), bold=True, color=G9, align=PP_ALIGN.CENTER, rtl=False)
try:
    resp = requests.get(
        'https://api.qrserver.com/v1/create-qr-code/'
        '?size=500x500&data=https%3A%2F%2Fdashboard-140293665526.me-west1.run.app'
        '&color=1b4332&bgcolor=ffffff&qzone=1', timeout=10)
    if resp.ok:
        qr_sz = qr_card_w - px(56)
        add_pic(sl, io.BytesIO(resp.content),
                qr_x + (qr_card_w-qr_sz)//2, qr_y+py(68), qr_sz, qr_sz)
        print('  QR downloaded')
except Exception as e:
    print(f'  QR skipped: {e}')
txb_lines(sl, ['שם משתמש: example@netivei.co.il', 'סיסמה: 123456'],
          qr_x+px(20), qr_y+qr_card_h-py(130), qr_card_w-px(30), py(120),
          size=fz(16), color=G9, align=PP_ALIGN.CENTER)

# Text — RIGHT side (first child in RTL)
tx = qr_x + qr_card_w + px(60)
tw2 = W - tx - px(PAD)
lbl(sl, 'מוזמנים להתרשם!', ty=60)   # uses full-width helper
txb_lines(sl, ['אי אפשר לצמצם', 'את מה שלא מודדים'],
          tx, py(140), tw2, py(450), size=fz(64), bold=True, color=G9)
txb(sl, 'CarbonTrack — נתיבי ישראל',
    tx, py(610), tw2, py(70), size=fz(28), bold=True, color=G5)
print('S15 Closing')


# ── Save ──────────────────────────────────────────────────────────────────────
out = BASE / 'CarbonTrack_Presentation_v2.pptx'
prs.save(str(out))
print(f'\nSaved: {out}')
